"""Genera archivos de prueba para cada tipo de documento soportado."""
from pathlib import Path

# 1. Word (.docx)
def create_docx(path):
    try:
        from docx import Document
    except ImportError:
        print(f'  (python-docx no instalado, se omite {path.name})')
        return
    doc = Document()
    doc.core_properties.title = 'Manual de pruebas'
    doc.core_properties.author = 'Jorge Aguayo'
    doc.core_properties.subject = 'QA'
    doc.core_properties.keywords = 'pdf, ninja, test'
    doc.add_heading('Capitulo 1', 1)
    doc.add_paragraph('Este es un parrafo de prueba.')
    table1 = doc.add_table(rows=2, cols=3)
    table1.cell(0, 0).text = 'A'
    table1.cell(0, 1).text = 'B'
    table1.cell(0, 2).text = 'C'
    table1.cell(1, 0).text = '1'
    table1.cell(1, 1).text = '2'
    table1.cell(1, 2).text = '3'
    table2 = doc.add_table(rows=3, cols=2)
    table2.cell(0, 0).text = 'X'
    doc.add_paragraph('Otro parrafo.')
    doc.save(str(path))
    print(f'OK {path.name}')


# 2. Excel (.xlsx)
def create_xlsx(path):
    import openpyxl
    wb = openpyxl.Workbook()
    wb.properties.title = 'Reporte Q3'
    wb.properties.creator = 'Jorge Aguayo'
    wb.properties.subject = 'Finanzas'
    wb.properties.keywords = 'reporte, finanzas'
    ws1 = wb.active
    ws1.title = 'Ventas'
    ws1['A1'] = 'Mes'
    ws1['B1'] = 'Monto'
    for i, m in enumerate(['Ene', 'Feb', 'Mar'], 2):
        ws1.cell(i, 1, m)
        ws1.cell(i, 2, 1000 * (i - 1))
    ws2 = wb.create_sheet('Costos')
    ws2['A1'] = 'Tipo'
    ws2['B1'] = 'Valor'
    for i, t in enumerate(['Fijo', 'Variable'], 2):
        ws2.cell(i, 1, t)
        ws2.cell(i, 2, 500 * i)
    ws3 = wb.create_sheet('Resumen', 0)
    wb.save(str(path))
    print(f'OK {path.name}')


# 3. PowerPoint (.pptx)
def create_pptx(path):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print(f'  (python-pptx no instalado, se omite {path.name})')
        return
    prs = Presentation()
    prs.core_properties.title = 'Presentacion de prueba'
    prs.core_properties.author = 'Jorge Aguayo'
    prs.core_properties.subject = 'Demo'
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Slide 1'
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Slide 2 con tabla'
    rows, cols = 2, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(5), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = 'Col1'
    table.cell(0, 1).text = 'Col2'
    table.cell(0, 2).text = 'Col3'
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Slide 3'
    prs.save(str(path))
    print(f'OK {path.name}')


# 4. Imagen con EXIF
def create_image(path):
    from PIL import Image, ImageDraw
    img = Image.new('RGB', (800, 600), color=(73, 109, 137))
    d = ImageDraw.Draw(img)
    d.text((10, 10), 'Test image', fill=(255, 255, 255))
    exif = img.getexif()
    exif[0x010F] = 'TestCam'
    exif[0x0110] = 'Model X'
    exif[0x0131] = 'PDF Ninja Test'
    img.save(str(path), exif=exif)
    print(f'OK {path.name}')


# 5. TXT
def create_txt(path):
    path.write_text('Linea 1\nLinea 2 con mas texto\nLinea 3\n', encoding='utf-8')
    print(f'OK {path.name}')


# 6. CSV
def create_csv(path):
    path.write_text('nombre,edad,ciudad\nAna,30,Santiago\nJuan,25,Valparaiso\nMaria,35,Concepcion\n', encoding='utf-8')
    print(f'OK {path.name}')


if __name__ == '__main__':
    create_docx(Path('test_doc.docx'))
    create_xlsx(Path('test_xls.xlsx'))
    create_pptx(Path('test_ppt.pptx'))
    create_image(Path('test_img.png'))
    create_txt(Path('test_text.txt'))
    create_csv(Path('test_data.csv'))
